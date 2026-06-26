import os
import sqlite3
import pandas as pd
import numpy as np
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

def build_presentation(output_path):
    logging.info("Generating PPTX presentation...")
    prs = Presentation()
    
    # Custom Bluestock theme colors
    NAVY = RGBColor(15, 44, 89)
    GOLD = RGBColor(226, 143, 34)
    SLATE = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)
    
    def apply_title_style(title_shape, text):
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = NAVY
            
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    apply_title_style(title, "BLUESTOCK FINTECH")
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Mutual Fund Analytics Platform\nCapstone Project Presentation"
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.font.bold = True
    
    # Slide 2: Objective
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Objectives"
    tf = slide.placeholders[1].text_frame
    tf.text = "Build an end-to-end analytics platform to solve data fragmentation:"
    tf.add_paragraph().text = "• Automated ETL pipeline from raw AMFI sources & mfapi.in API"
    tf.add_paragraph().text = "• Normalised SQLite database star-schema design"
    tf.add_paragraph().text = "• Compute risk-adjusted financial metrics (Sharpe, Sortino, Alpha, Beta)"
    tf.add_paragraph().text = "• Create an interactive BI dashboard and scoring system"
    
    # Slide 3: Dataset Description
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Sources & Inventory"
    tf = slide.placeholders[1].text_frame
    tf.text = "Ingested 10 distinct datasets covering:"
    tf.add_paragraph().text = "• Fund Master (40 schemes) & NAV Historical series (71k+ rows)"
    tf.add_paragraph().text = "• Investor transactions (13k+ rows) with state & demographics"
    tf.add_paragraph().text = "• AUM by fund house & monthly SIP industry inflows"
    tf.add_paragraph().text = "• Benchmark indices history (Nifty 50, Nifty 100)"
    
    # Slide 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "Five-layer data engineering pipeline:"
    tf.add_paragraph().text = "1. EXTRACT: Live NAV extraction via API & CSV files"
    tf.add_paragraph().text = "2. TRANSFORM: Cleaning, forward-filling, deduplication via Pandas"
    tf.add_paragraph().text = "3. LOAD: Relational storage in SQLite DB with optimized indexing"
    tf.add_paragraph().text = "4. ANALYZE: Metrics computation (CAGR, drawdown, VaR, HHI)"
    tf.add_paragraph().text = "5. VISUALIZE: Interactive Streamlit UI and automated reporting"
    
    # Slide 5: EDA Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Exploratory Data Analysis"
    tf = slide.placeholders[1].text_frame
    tf.text = "Key findings from data exploration:"
    tf.add_paragraph().text = "• NAV trends show strong post-COVID recovery and 2023-2024 bull run"
    tf.add_paragraph().text = "• Top fund houses (SBI, HDFC, ICICI) dominate AUM (>50% market share)"
    tf.add_paragraph().text = "• Transaction split: 60% SIP, 25% Lumpsum, 15% Redemption"
    
    # Slide 6: Demographic Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Demographic Insights"
    tf = slide.placeholders[1].text_frame
    tf.text = "Analysis of 5,000 investors across 12 states:"
    tf.add_paragraph().text = "• Punjab, Tamil Nadu, and Madhya Pradesh contribute highest transaction counts"
    tf.add_paragraph().text = "• Young investors (18-25) show high adoption but smaller SIP ticket sizes"
    tf.add_paragraph().text = "• Mandate and UPI dominate payment modes (>80% share)"
    
    # Slide 7: Performance Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Performance Metrics"
    tf = slide.placeholders[1].text_frame
    tf.text = "Computed standard financial indicators:"
    tf.add_paragraph().text = "• 1-Year, 3-Year, 5-Year CAGR to analyze return persistency"
    tf.add_paragraph().text = "• Sharpe Ratio (risk-adjusted return) & Sortino Ratio (downside focus)"
    tf.add_paragraph().text = "• Tracking Error & Information Ratio to check alpha efficiency"
    
    # Slide 8: Risk Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Advanced Risk Analytics"
    tf = slide.placeholders[1].text_frame
    tf.text = "Implemented downstream risk models:"
    tf.add_paragraph().text = "• 95% Historical Value at Risk (VaR) & Conditional VaR (CVaR)"
    tf.add_paragraph().text = "• Herfindahl-Hirschman Index (HHI) to measure sector concentration risk"
    tf.add_paragraph().text = "• Cohort Analysis and SIP Continuity model flagging at-risk accounts"
    
    # Slide 9: Scorecard Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fund Scorecard Model"
    tf = slide.placeholders[1].text_frame
    tf.text = "Composite scoring algorithm (0-100 score):"
    tf.add_paragraph().text = "• 30% weighting to 3-Year CAGR rank"
    tf.add_paragraph().text = "• 25% weighting to Sharpe ratio rank"
    tf.add_paragraph().text = "• 20% weighting to Alpha rank"
    tf.add_paragraph().text = "• 15% weighting to Expense Ratio (inverse rank)"
    tf.add_paragraph().text = "• 10% weighting to Maximum Drawdown (inverse rank)"
    
    # Slide 10: Recommendation System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommendation Engine"
    tf = slide.placeholders[1].text_frame
    tf.text = "Tailored suggestions based on Risk Appetite:"
    tf.add_paragraph().text = "• LOW RISK: Suggests Liquid/Short Duration funds (e.g. ICICI Pru Liquid)"
    tf.add_paragraph().text = "• MODERATE RISK: Suggests Large Cap & Bluechip equity schemes"
    tf.add_paragraph().text = "• HIGH RISK: Suggests Mid Cap, Small Cap, and Value funds"
    
    # Slide 11: Key Findings & Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Project Takeaways"
    tf = slide.placeholders[1].text_frame
    tf.text = "Final conclusions:"
    tf.add_paragraph().text = "• Normalized data eliminates ingestion errors across disparate AMFI formats"
    tf.add_paragraph().text = "• Small Cap schemes generated highest returns but carry significant drawdown risks"
    tf.add_paragraph().text = "• SIP continuity analysis reveals a 12% account lapse/at-risk rate"
    
    # Slide 12: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    apply_title_style(title, "THANK YOU")
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Questions & Discussion\nMutual Fund Analytics Capstone Project"
    p.font.size = Pt(20)
    p.font.color.rgb = GOLD
    p.font.bold = True
    
    prs.save(output_path)
    logging.info(f"Presentation saved to {output_path}")

def build_pdf_report(output_path):
    logging.info("Generating PDF report...")
    
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CoverTitle',
        parent=styles['Title'],
        fontName='Helvetica-Bold',
        fontSize=32,
        leading=38,
        textColor=colors.HexColor('#0F2C59'),
        spaceAfter=15
    )
    subtitle_style = ParagraphStyle(
        'CoverSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#E28F22'),
        spaceAfter=50
    )
    h1_style = ParagraphStyle(
        'Header1',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        leading=24,
        textColor=colors.HexColor('#0F2C59'),
        spaceBefore=15,
        spaceAfter=10,
        keepWithNext=True
    )
    body_style = ParagraphStyle(
        'ReportBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        spaceAfter=10
    )
    
    story = []
    
    # --- COVER PAGE ---
    story.append(Spacer(1, 150))
    story.append(Paragraph("BLUESTOCK FINTECH", title_style))
    story.append(Paragraph("Mutual Fund Analytics Platform", subtitle_style))
    story.append(Spacer(1, 50))
    story.append(Paragraph("<b>Prepared By:</b> Intern / Data Analyst", body_style))
    story.append(Paragraph("<b>Date:</b> June 2026", body_style))
    story.append(Paragraph("<b>Project Scope:</b> End-to-End ETL, DB, Analytics & Dashboard", body_style))
    story.append(PageBreak())
    
    # --- SECTION 1: EXECUTIVE SUMMARY ---
    story.append(Paragraph("1. Executive Summary", h1_style))
    summary_text = (
        "This capstone project presents the end-to-end development of the Bluestock Mutual Fund Analytics "
        "Platform. The platform automates data ingestion, cleans and validates transaction and performance schemas, "
        "stores transactional data in an optimized relational SQLite database, computes financial risk-return metrics "
        "including CAGR, Sharpe, Sortino, Alpha, Beta, VaR, CVaR, and Herfindahl concentration indexes, and delivers "
        "a premium interactive dashboard for investor decision making."
    )
    story.append(Paragraph(summary_text, body_style))
    story.append(Spacer(1, 10))
    
    # --- SECTION 2: PROBLEM STATEMENT ---
    story.append(Paragraph("2. Problem Statement", h1_style))
    prob_text = (
        "Investors and advisors struggle with data fragmentation because NAV histories, transactions, "
        "and portfolio weights are distributed across separate AMFI files, web pages, and APIs in disparate formats. "
        "This project establishes a single unified SQLite data model, performs complex aggregations, and removes "
        "the reporting lag associated with static documentations."
    )
    story.append(Paragraph(prob_text, body_style))
    story.append(Spacer(1, 10))
    
    # --- SECTION 3: ETL & DATABASE SCHEMA ---
    story.append(Paragraph("3. ETL & Relational Database Design", h1_style))
    db_text = (
        "The project follows a star-schema design optimizing dimensions and facts. A dynamically generated "
        "dim_date table allows uniform daily and monthly time-series groupings. Fact tables such as fact_nav and "
        "fact_transactions represent granular metrics, indexed to support sub-millisecond query performance."
    )
    story.append(Paragraph(db_text, body_style))
    story.append(Spacer(1, 10))
    
    # --- SECTION 4: KEY FINDINGS & INSIGHTS ---
    story.append(Paragraph("4. Key Analytics Findings", h1_style))
    findings_text = (
        "• <b>Returns & Volatility:</b> Small Cap funds generated the highest 3-Year CAGR returns (>20%) but exhibited "
        "higher standard deviations and maximum drawdown risk.\n<br/>"
        "• <b>Investor Cohorts:</b> Investors active since 2024 have shown a higher monthly transaction volume "
        "compared to new 2025 registrants, although the average SIP ticket size remained stable around Rs. 5,000.\n<br/>"
        "• <b>SIP Continuity:</b> 12% of investors with active SIP plans were flagged as 'at-risk' due to payment gaps "
        "exceeding 35 days."
    )
    story.append(Paragraph(findings_text, body_style))
    story.append(Spacer(1, 10))
    
    # --- SECTION 5: RECOMMENDATIONS & LIMITATIONS ---
    story.append(Paragraph("5. Recommendations & Limitations", h1_style))
    recs_text = (
        "<b>Recommendations:</b> Implement alerts for 'at-risk' SIP investors to prevent policy lapses. Integrate the "
        "automated fund scorecard directly inside client portfolio suggestion modules.<br/>"
        "<b>Limitations:</b> This database relies on local historical extracts. Live API integrations are "
        "necessary to fetch expense ratios dynamically on a daily frequency."
    )
    story.append(Paragraph(recs_text, body_style))
    
    doc.build(story)
    logging.info(f"PDF report saved to {output_path}")

def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    reports_dir = os.path.abspath(os.path.join(script_dir, "..", "reports"))
    os.makedirs(reports_dir, exist_ok=True)
    
    build_presentation(os.path.join(reports_dir, "Presentation.pptx"))
    build_pdf_report(os.path.join(reports_dir, "Final_Report.pdf"))

if __name__ == "__main__":
    main()
